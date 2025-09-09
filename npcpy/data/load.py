import fitz
import pandas as pd
import json
import io
from PIL import Image
import numpy as np
from typing import Optional
import os

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None

def load_csv(file_path):
    df = pd.read_csv(file_path)
    return df

def load_json(file_path):
    with open(file_path, "r", encoding='utf-8') as f:
        data = json.load(f)
    return data

def load_txt(file_path):
    with open(file_path, "r", encoding='utf-8') as f:
        text = f.read()
    return text

def load_excel(file_path):
    df = pd.read_excel(file_path)
    return df

def load_image(file_path):
    img = Image.open(file_path)
    img_array = np.array(img)
    df = pd.DataFrame(
        {
            "image_array": [img_array.tobytes()],
            "shape": [img_array.shape],
            "dtype": [img_array.dtype.str],
        }
    )
    return df

def load_pdf(file_path):
    pdf_document = fitz.open(file_path)
    full_text = ""
    for page in pdf_document:
        full_text += page.get_text() + "\n"
    return full_text

def load_docx(file_path):
    if Document is None:
        raise ImportError("Please install python-docx to load .docx files.")
    doc = Document(file_path)
    full_text = "\n".join([para.text for para in doc.paragraphs])
    return full_text

def load_pptx(file_path):
    if Presentation is None:
        raise ImportError("Please install python-pptx to load .pptx files.")
    prs = Presentation(file_path)
    full_text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                full_text += shape.text + "\n"
    return full_text

def load_html(file_path):
    if BeautifulSoup is None:
        raise ImportError("Please install beautifulsoup4 to load .html files.")
    with open(file_path, 'r', encoding='utf-8') as f:
        soup = BeautifulSoup(f, 'html.parser')
    return soup.get_text(separator='\n', strip=True)

extension_map = {
    "PNG": "images",
    "JPG": "images",
    "JPEG": "images",
    "GIF": "images",
    "SVG": "images",
    "MP4": "videos",
    "AVI": "videos",
    "MOV": "videos",
    "WMV": "videos",
    "MPG": "videos",
    "MPEG": "videos",
    "DOCX": "documents",
    "PPTX": "documents",
    "PDF": "documents",
    "XLSX": "documents",
    "TXT": "documents",
    "CSV": "documents",
    "MD": "documents",
    "HTML": "documents",
    "HTM": "documents",
    "ZIP": "archives",
    "RAR": "archives",
    "7Z": "archives",
    "TAR": "archives",
    "GZ": "archives",
}

def load_file_contents(file_path, chunk_size=None):
    file_ext = os.path.splitext(file_path)[1].upper().lstrip('.')
    full_content = ""
    if not isinstance(chunk_size, int):
        chunk_size=250
    try:
        if file_ext == 'PDF':
            full_content = load_pdf(file_path)
        elif file_ext == 'DOCX':
            full_content = load_docx(file_path)
        elif file_ext == 'PPTX':
            full_content = load_pptx(file_path)
        elif file_ext in ['HTML', 'HTM']:
            full_content = load_html(file_path)
        elif file_ext == 'CSV':
            df = load_csv(file_path)
            full_content = df.to_string()
        elif file_ext in ['XLS', 'XLSX']:
            df = load_excel(file_path)
            full_content = df.to_string()
        elif file_ext in ['TXT', 'MD', 'PY', 'JSX', 'TSX', 'TS', 'JS', 'JSON', 'SQL', 'NPC', 'JINX', 'LINE', 'YAML']:
            full_content = load_txt(file_path)
        elif file_ext == 'JSON':
            data = load_json(file_path)
            full_content = json.dumps(data, indent=2)
        else:
            return [f"Unsupported file format for content loading: {file_ext}"]
        
        if not full_content:
            return []

        chunks = []
        for i in range(0, len(full_content), chunk_size):
            chunk = full_content[i:i+chunk_size].strip()
            if chunk:
                chunks.append(chunk)
        return chunks
            
    except Exception as e:
        return [f"Error loading file {file_path}: {str(e)}"]
